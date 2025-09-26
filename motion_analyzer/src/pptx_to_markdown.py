#!/usr/bin/env python3
"""
PPTX to Markdown Converter
Converts PowerPoint presentations to detailed Markdown format with all content preserved.
"""

import os
import sys
import json
import base64
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io
except ImportError:
    print("Installing required packages...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "Pillow"])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io


class PPTXToMarkdownConverter:
    """Convert PowerPoint presentations to detailed Markdown format."""

    def __init__(self, pptx_path: str, output_dir: str = None):
        """Initialize converter with PPTX file path."""
        self.pptx_path = Path(pptx_path)
        self.presentation = Presentation(pptx_path)

        # Set up output directory
        if output_dir:
            self.output_dir = Path(output_dir)
        else:
            self.output_dir = self.pptx_path.parent / "converted"

        self.output_dir.mkdir(exist_ok=True)
        self.images_dir = self.output_dir / "images"
        self.images_dir.mkdir(exist_ok=True)

        # Store extracted content
        self.slides_content = []
        self.metadata = {}

    def extract_metadata(self) -> Dict:
        """Extract presentation metadata."""
        core_props = self.presentation.core_properties

        metadata = {
            "title": core_props.title or self.pptx_path.stem,
            "author": core_props.author or "Unknown",
            "subject": core_props.subject or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "slide_count": len(self.presentation.slides),
            "file_name": self.pptx_path.name,
            "converted_date": datetime.now().isoformat()
        }

        return metadata

    def extract_text_from_shape(self, shape) -> str:
        """Extract text from a shape."""
        if not shape.has_text_frame:
            return ""

        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            para_text = []
            for run in paragraph.runs:
                para_text.append(run.text)

            full_para = "".join(para_text).strip()
            if full_para:
                # Add proper markdown formatting based on paragraph level
                if paragraph.level == 0:
                    text_parts.append(full_para)
                else:
                    # Add indentation for nested items
                    indent = "  " * paragraph.level
                    text_parts.append(f"{indent}- {full_para}")

        return "\n".join(text_parts)

    def extract_table_from_shape(self, shape) -> str:
        """Extract table from a shape and format as markdown."""
        if not shape.has_table:
            return ""

        table = shape.table
        markdown_lines = []

        # Extract headers (first row)
        headers = []
        for cell in table.rows[0].cells:
            headers.append(cell.text.strip())

        # Create markdown table header
        markdown_lines.append("| " + " | ".join(headers) + " |")
        markdown_lines.append("|" + "|".join(["---" for _ in headers]) + "|")

        # Extract data rows
        for row_idx in range(1, len(table.rows)):
            row_data = []
            for cell in table.rows[row_idx].cells:
                row_data.append(cell.text.strip())
            markdown_lines.append("| " + " | ".join(row_data) + " |")

        return "\n".join(markdown_lines)

    def save_image_from_shape(self, shape, slide_num: int, img_num: int) -> Optional[str]:
        """Save image from shape and return relative path."""
        try:
            image = shape.image
            image_bytes = image.blob

            # Determine file extension
            ext = image.ext
            if not ext:
                ext = 'png'

            # Save image
            image_filename = f"slide_{slide_num:03d}_img_{img_num:03d}.{ext}"
            image_path = self.images_dir / image_filename

            with open(image_path, 'wb') as f:
                f.write(image_bytes)

            return f"images/{image_filename}"
        except Exception as e:
            print(f"Error saving image: {e}")
            return None

    def process_slide(self, slide, slide_num: int) -> Dict:
        """Process a single slide and extract all content."""
        slide_content = {
            "slide_number": slide_num,
            "layout": slide.slide_layout.name if slide.slide_layout else "Custom",
            "title": "",
            "content": [],
            "notes": "",
            "images": [],
            "tables": [],
            "shapes_info": []
        }

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_content["notes"] = slide.notes_slide.notes_text_frame.text

        # Process all shapes
        img_counter = 1
        for shape in slide.shapes:
            shape_info = {
                "type": shape.shape_type,
                "name": shape.name,
                "position": f"({shape.left}, {shape.top})",
                "size": f"({shape.width}, {shape.height})"
            }

            # Extract text
            if shape.has_text_frame:
                text = self.extract_text_from_shape(shape)
                if text:
                    # Check if it's likely a title
                    if shape.name.lower().startswith("title") or shape == slide.shapes.title:
                        slide_content["title"] = text
                    else:
                        slide_content["content"].append({
                            "type": "text",
                            "content": text,
                            "shape_name": shape.name
                        })

            # Extract tables
            if shape.has_table:
                table_markdown = self.extract_table_from_shape(shape)
                if table_markdown:
                    slide_content["tables"].append({
                        "content": table_markdown,
                        "shape_name": shape.name
                    })

            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_path = self.save_image_from_shape(shape, slide_num, img_counter)
                if image_path:
                    slide_content["images"].append({
                        "path": image_path,
                        "shape_name": shape.name,
                        "alt_text": f"Image {img_counter} from slide {slide_num}"
                    })
                    img_counter += 1

            slide_content["shapes_info"].append(shape_info)

        return slide_content

    def convert_to_markdown(self) -> str:
        """Convert entire presentation to markdown."""
        # Extract metadata
        self.metadata = self.extract_metadata()

        # Process all slides
        for idx, slide in enumerate(self.presentation.slides, 1):
            slide_content = self.process_slide(slide, idx)
            self.slides_content.append(slide_content)

        # Generate markdown
        markdown_lines = []

        # Add metadata header
        markdown_lines.append(f"# {self.metadata['title']}")
        markdown_lines.append("")
        markdown_lines.append("## Metadata")
        markdown_lines.append(f"- **Author**: {self.metadata['author']}")
        markdown_lines.append(f"- **Created**: {self.metadata['created']}")
        markdown_lines.append(f"- **Modified**: {self.metadata['modified']}")
        markdown_lines.append(f"- **Total Slides**: {self.metadata['slide_count']}")
        markdown_lines.append(f"- **File**: {self.metadata['file_name']}")
        markdown_lines.append(f"- **Converted**: {self.metadata['converted_date']}")
        markdown_lines.append("")
        markdown_lines.append("---")
        markdown_lines.append("")

        # Add table of contents
        markdown_lines.append("## Table of Contents")
        markdown_lines.append("")
        for slide in self.slides_content:
            title = slide['title'] or f"Slide {slide['slide_number']}"
            markdown_lines.append(f"- [Slide {slide['slide_number']}: {title}](#slide-{slide['slide_number']})")
        markdown_lines.append("")
        markdown_lines.append("---")
        markdown_lines.append("")

        # Add slides content
        for slide in self.slides_content:
            # Slide header
            markdown_lines.append(f"## Slide {slide['slide_number']}")
            markdown_lines.append("")

            # Title
            if slide['title']:
                markdown_lines.append(f"### {slide['title']}")
                markdown_lines.append("")

            # Layout info
            markdown_lines.append(f"**Layout**: {slide['layout']}")
            markdown_lines.append("")

            # Content sections
            if slide['content']:
                markdown_lines.append("### Content")
                markdown_lines.append("")
                for content_item in slide['content']:
                    markdown_lines.append(content_item['content'])
                    markdown_lines.append("")

            # Tables
            if slide['tables']:
                markdown_lines.append("### Tables")
                markdown_lines.append("")
                for table in slide['tables']:
                    markdown_lines.append(table['content'])
                    markdown_lines.append("")

            # Images
            if slide['images']:
                markdown_lines.append("### Images")
                markdown_lines.append("")
                for img in slide['images']:
                    markdown_lines.append(f"![{img['alt_text']}]({img['path']})")
                    markdown_lines.append("")

            # Notes
            if slide['notes']:
                markdown_lines.append("### Speaker Notes")
                markdown_lines.append("")
                markdown_lines.append(f"_{slide['notes']}_")
                markdown_lines.append("")

            # Separator
            markdown_lines.append("---")
            markdown_lines.append("")

        return "\n".join(markdown_lines)

    def save_markdown(self, markdown_content: str) -> str:
        """Save markdown content to file."""
        output_file = self.output_dir / f"{self.pptx_path.stem}.md"

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(markdown_content)

        return str(output_file)

    def save_json_data(self) -> str:
        """Save structured data as JSON for programmatic access."""
        json_file = self.output_dir / f"{self.pptx_path.stem}.json"

        data = {
            "metadata": self.metadata,
            "slides": self.slides_content
        }

        with open(json_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, default=str)

        return str(json_file)

    def convert(self) -> Dict[str, str]:
        """Main conversion method."""
        print(f"Converting {self.pptx_path.name}...")

        # Convert to markdown
        markdown_content = self.convert_to_markdown()

        # Save files
        md_file = self.save_markdown(markdown_content)
        json_file = self.save_json_data()

        print(f"✅ Conversion complete!")
        print(f"📄 Markdown: {md_file}")
        print(f"📊 JSON: {json_file}")
        print(f"🖼️  Images: {self.images_dir}")

        return {
            "markdown_file": md_file,
            "json_file": json_file,
            "images_dir": str(self.images_dir)
        }


def main():
    """Main function for command-line usage."""
    if len(sys.argv) < 2:
        print("Usage: python pptx_to_markdown.py <path_to_pptx_file> [output_directory]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(pptx_file):
        print(f"Error: File '{pptx_file}' not found.")
        sys.exit(1)

    converter = PPTXToMarkdownConverter(pptx_file, output_dir)
    result = converter.convert()

    return result


if __name__ == "__main__":
    main()